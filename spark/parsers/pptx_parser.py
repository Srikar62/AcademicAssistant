"""
PPTX parser — extracts text from PowerPoint presentations using python-pptx.

Handles run-merging properly (the reason this exists instead of raw XML parsing).
Returns one ParsedSection per slide, including the slide title as metadata.
"""

import logging
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches
import io

from spark.parsers.base import BaseParser, ParsedSection

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """Extract text from PowerPoint slides with proper run-merging."""

    def parse(self, file_bytes: bytes, filename: str) -> List[ParsedSection]:
        """
        Parse a PPTX from raw bytes.

        Args:
            file_bytes: Raw PPTX content.
            filename: Original filename for logging.

        Returns:
            One ParsedSection per slide with non-empty text.

        Raises:
            ValueError: If the PPTX is corrupted or unreadable.
        """
        try:
            prs = Presentation(io.BytesIO(file_bytes))
        except Exception as exc:
            raise ValueError(
                f"Failed to open PPTX '{filename}': {exc}"
            ) from exc

        sections: List[ParsedSection] = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_title = self._extract_title(slide)
            slide_text = self._extract_slide_text(slide)
            slide_tables = self._extract_tables(slide)

            if not slide_text.strip():
                logger.debug(
                    "Slide %d of '%s' has no text content.",
                    slide_idx + 1,
                    filename,
                )
                continue

            metadata = {
                "slide_number": slide_idx + 1,
                "slide_title": slide_title or "",
                "total_slides": len(prs.slides),
            }
            if slide_tables:
                metadata["tables"] = slide_tables

            sections.append(
                ParsedSection(
                    label=f"Slide {slide_idx + 1}",
                    index=slide_idx,
                    text=slide_text.strip(),
                    metadata=metadata,
                )
            )

        if not sections:
            raise ValueError(
                f"PPTX '{filename}' yielded no extractable text "
                f"({len(prs.slides)} slides scanned)."
            )

        logger.info(
            "Parsed '%s': %d slides, %d with text.",
            filename,
            len(prs.slides),
            len(sections),
        )
        return sections

    @staticmethod
    def _extract_title(slide) -> Optional[str]:
        """
        Extract the slide title from the title placeholder.
        Returns None if no title placeholder exists.
        """
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text_frame.text.strip()
        return None

    @staticmethod
    def _extract_slide_text(slide) -> str:
        """
        Extract all text from a slide, merging runs within paragraphs.

        python-pptx handles run-merging correctly — each paragraph's runs
        are concatenated, and paragraphs are joined with newlines.  This
        avoids the raw-XML fragmentation problem (splitting mid-word on
        <a:t> boundaries) that the client-side plan flagged.

        Tables are rendered as Markdown tables for readable text output.
        Structured table data is collected separately via _extract_tables().
        """
        text_parts: List[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Merge all runs in this paragraph into a single string
                    para_text = "".join(run.text for run in paragraph.runs).strip()
                    if para_text:
                        text_parts.append(para_text)

            # Render tables as Markdown for readable embedding text
            if shape.has_table:
                md_table = PPTXParser._table_to_markdown(shape.table)
                if md_table:
                    text_parts.append(md_table)

        return "\n".join(text_parts)

    @staticmethod
    def _extract_tables(slide) -> List[List[List[str]]]:
        """
        Extract structured table data from a slide.

        Returns:
            List of tables, where each table is a list of rows,
            and each row is a list of cell strings.
        """
        tables: List[List[List[str]]] = []
        for shape in slide.shapes:
            if shape.has_table:
                table_data: List[List[str]] = []
                for row in shape.table.rows:
                    table_data.append([cell.text.strip() for cell in row.cells])
                if table_data and any(any(c for c in row) for row in table_data):
                    tables.append(table_data)
        return tables

    @staticmethod
    def _table_to_markdown(table) -> str:
        """
        Convert a python-pptx Table object to a Markdown table string.

        The first row is treated as the header.  Returns an empty string
        if the table has no meaningful content.
        """
        rows: List[List[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])

        # Filter out completely empty tables
        if not rows or not any(any(c for c in row) for row in rows):
            return ""

        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
        body = "\n".join(
            "| " + " | ".join(row) + " |" for row in rows[1:]
        )

        parts = [header, separator]
        if body:
            parts.append(body)
        return "\n".join(parts)

