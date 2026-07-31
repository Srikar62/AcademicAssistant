"""
Integration tests for the Spark processing pipeline.

Tests are organized by component (parsers, chunker, embedder, pipeline)
and use real file fixtures where practical.  External services (MinIO,
Qdrant, Kafka) are mocked so tests run without Docker.
"""

import io
import os
import json
import pytest
from unittest.mock import MagicMock, patch, PropertyMock
from dataclasses import asdict

# ─── Test fixtures: minimal real files ─────────────────────────


def _make_pdf_bytes() -> bytes:
    """Create a minimal valid PDF with text using PyMuPDF."""
    import fitz

    doc = fitz.open()
    for i in range(3):
        page = doc.new_page()
        text = f"This is page {i + 1}. It contains important academic content about machine learning algorithms and neural network architectures."
        page.insert_text((72, 72), text, fontsize=12)
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def _make_pptx_bytes() -> bytes:
    """Create a minimal valid PPTX with text using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # Slide 1 — title slide
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Introduction to Machine Learning"
    slide.placeholders[1].text = "A comprehensive overview of ML concepts"

    # Slide 2 — content
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Supervised Learning"
    slide.placeholders[1].text = (
        "Supervised learning uses labeled data to train models. "
        "Common algorithms include linear regression, decision trees, "
        "and support vector machines."
    )

    # Slide 3 — content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Unsupervised Learning"
    slide.placeholders[1].text = (
        "Unsupervised learning discovers patterns in unlabeled data. "
        "Clustering and dimensionality reduction are key techniques."
    )

    # Slide 4 — very short (to test merging)
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Summary"
    slide.placeholders[1].text = "Key takeaways."

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_text_bytes() -> bytes:
    """Create sample text content."""
    paragraphs = [
        "Machine learning is a branch of artificial intelligence that focuses on building systems that learn from data.",
        "Deep learning is a subset of machine learning that uses neural networks with many layers to learn complex patterns.",
        "Natural language processing enables computers to understand, interpret, and generate human language.",
        "Computer vision allows machines to interpret and make decisions based on visual data from the world.",
        "Reinforcement learning is a type of machine learning where an agent learns to make decisions by taking actions.",
        "Transfer learning leverages knowledge from one task to improve performance on a related task.",
    ]
    return "\n\n".join(paragraphs).encode("utf-8")


# ═══════════════════════════════════════════════════════════════
#  Parser Tests
# ═══════════════════════════════════════════════════════════════


class TestPDFParser:
    """Tests for the PDF parser."""

    def test_parse_valid_pdf(self):
        from spark.parsers.pdf_parser import PDFParser

        parser = PDFParser()
        sections = parser.parse(_make_pdf_bytes(), "test.pdf")

        assert len(sections) == 3
        assert sections[0].label == "Page 1"
        assert sections[2].label == "Page 3"
        assert "machine learning" in sections[0].text.lower()
        assert sections[0].metadata["page_number"] == 1
        assert sections[0].metadata["total_pages"] == 3

    def test_parse_corrupted_pdf_raises(self):
        from spark.parsers.pdf_parser import PDFParser

        parser = PDFParser()
        with pytest.raises(ValueError, match="Failed to open PDF"):
            parser.parse(b"not a pdf", "corrupted.pdf")


class TestPPTXParser:
    """Tests for the PPTX parser."""

    def test_parse_valid_pptx(self):
        from spark.parsers.pptx_parser import PPTXParser

        parser = PPTXParser()
        sections = parser.parse(_make_pptx_bytes(), "slides.pptx")

        assert len(sections) >= 3
        # Check that slide titles are extracted
        has_ml_title = any(
            "machine learning" in s.metadata.get("slide_title", "").lower()
            for s in sections
        )
        assert has_ml_title

    def test_parse_corrupted_pptx_raises(self):
        from spark.parsers.pptx_parser import PPTXParser

        parser = PPTXParser()
        with pytest.raises(ValueError, match="Failed to open PPTX"):
            parser.parse(b"not a pptx", "corrupted.pptx")


class TestTextParser:
    """Tests for the text/markdown parser."""

    def test_parse_valid_text(self):
        from spark.parsers.text_parser import TextParser

        parser = TextParser()
        sections = parser.parse(_make_text_bytes(), "notes.txt")

        # 6 paragraphs → should be split into sections
        assert len(sections) == 6
        assert "machine learning" in sections[0].text.lower()

    def test_parse_short_text_single_section(self):
        from spark.parsers.text_parser import TextParser

        parser = TextParser()
        short = b"Just a short note.\n\nWith two paragraphs."
        sections = parser.parse(short, "short.txt")

        # < 5 paragraphs → single section
        assert len(sections) == 1
        assert sections[0].label == "Full Document"

    def test_parse_empty_text_raises(self):
        from spark.parsers.text_parser import TextParser

        parser = TextParser()
        with pytest.raises(ValueError, match="empty"):
            parser.parse(b"", "empty.txt")


class TestParserDispatcher:
    """Tests for the get_parser factory."""

    def test_get_pdf_parser(self):
        from spark.parsers.base import get_parser
        from spark.parsers.pdf_parser import PDFParser

        assert isinstance(get_parser(".pdf"), PDFParser)

    def test_get_pptx_parser(self):
        from spark.parsers.base import get_parser
        from spark.parsers.pptx_parser import PPTXParser

        assert isinstance(get_parser(".pptx"), PPTXParser)

    def test_get_txt_parser(self):
        from spark.parsers.base import get_parser
        from spark.parsers.text_parser import TextParser

        assert isinstance(get_parser(".txt"), TextParser)

    def test_get_md_parser(self):
        from spark.parsers.base import get_parser
        from spark.parsers.text_parser import TextParser

        assert isinstance(get_parser(".md"), TextParser)

    def test_unsupported_type_raises(self):
        from spark.parsers.base import get_parser

        with pytest.raises(ValueError, match="Unsupported"):
            get_parser(".exe")


# ═══════════════════════════════════════════════════════════════
#  Chunker Tests
# ═══════════════════════════════════════════════════════════════


class TestChunker:
    """Tests for the sentence-aware chunker."""

    def test_chunk_basic(self):
        from spark.parsers.base import ParsedSection
        from spark.chunker import chunk_sections

        sections = [
            ParsedSection(
                label="Page 1",
                index=0,
                text=(
                    "Machine learning is powerful. It learns from data. "
                    "Deep learning uses neural networks. These networks have many layers. "
                    "Training requires lots of data. GPUs accelerate training. "
                    "Models can overfit. Regularization prevents overfitting."
                ),
            ),
        ]

        chunks = chunk_sections(sections, max_tokens=20, overlap_fraction=0.15)

        assert len(chunks) >= 1
        # Each chunk should not wildly exceed max_tokens
        for chunk in chunks:
            assert chunk.token_count <= 30  # some slack for sentence boundaries

    def test_chunk_preserves_source_labels(self):
        from spark.parsers.base import ParsedSection
        from spark.chunker import chunk_sections

        sections = [
            ParsedSection(label="Page 1", index=0, text="First page content here."),
            ParsedSection(label="Page 2", index=1, text="Second page content here."),
        ]

        chunks = chunk_sections(sections, max_tokens=100)
        assert len(chunks) >= 1
        assert "Page" in chunks[0].source_label

    def test_chunk_overlap(self):
        from spark.parsers.base import ParsedSection
        from spark.chunker import chunk_sections

        # Create enough text to force multiple chunks
        long_text = ". ".join(
            [f"Sentence number {i} with enough words to count as tokens"
             for i in range(50)]
        ) + "."

        sections = [
            ParsedSection(label="Page 1", index=0, text=long_text),
        ]

        chunks = chunk_sections(sections, max_tokens=30, overlap_fraction=0.2)

        if len(chunks) >= 2:
            # Check that consecutive chunks share some text (overlap)
            words_0 = set(chunks[0].text.split())
            words_1 = set(chunks[1].text.split())
            overlap = words_0 & words_1
            assert len(overlap) > 0, "Expected overlap between consecutive chunks"

    def test_chunk_empty_sections(self):
        from spark.chunker import chunk_sections

        chunks = chunk_sections([])
        assert chunks == []


class TestPPTXChunker:
    """Tests for PPTX-specific chunking (slide merging)."""

    def test_merges_short_slides(self):
        from spark.parsers.base import ParsedSection
        from spark.chunker import chunk_pptx_sections

        sections = [
            ParsedSection(
                label="Slide 1", index=0,
                text="Introduction to the course.",  # very short
            ),
            ParsedSection(
                label="Slide 2", index=1,
                text="Overview of the syllabus.",  # very short
            ),
            ParsedSection(
                label="Slide 3", index=2,
                text=(
                    "Machine learning is a field that involves algorithms "
                    "that learn from data and improve over time without being "
                    "explicitly programmed for specific tasks."
                ),
            ),
        ]

        chunks = chunk_pptx_sections(
            sections, max_tokens=100, min_slide_tokens=50
        )

        assert len(chunks) >= 1
        # Short slides should have been merged
        merged_labels = [c.source_label for c in chunks]
        has_merged = any("–" in label for label in merged_labels)
        # Either merged or all fit in one chunk
        assert has_merged or len(chunks) <= 2


# ═══════════════════════════════════════════════════════════════
#  Embedder Tests (model mocked to avoid heavy downloads)
# ═══════════════════════════════════════════════════════════════


class TestEmbedder:
    """Tests for the embedding module."""

    def test_embed_chunks_batch(self):
        import numpy as np
        from spark.embedder import embed_chunks_batch

        # Mock the sentence-transformers model
        mock_model = MagicMock()
        mock_model.encode.return_value = np.random.rand(3, 384).astype(np.float32)

        with patch.dict("spark.embedder._MODEL_CACHE", {"test-model": mock_model}):
            chunks = [
                {"text": "Hello world"},
                {"text": "Machine learning is great"},
                {"text": "Deep learning uses neural networks"},
            ]
            result = embed_chunks_batch(chunks, model_name="test-model")

            assert len(result) == 3
            for chunk in result:
                assert "embedding" in chunk
                assert len(chunk["embedding"]) == 384

    def test_embed_empty_list(self):
        from spark.embedder import embed_chunks_batch

        result = embed_chunks_batch([], model_name="test-model")
        assert result == []


# ═══════════════════════════════════════════════════════════════
#  Qdrant Writer Tests (mocked client)
# ═══════════════════════════════════════════════════════════════


class TestQdrantWriter:
    """Tests for the Qdrant writer."""

    def test_ensure_collection_creates_when_missing(self):
        from spark.qdrant_writer import QdrantWriter

        mock_client = MagicMock()
        mock_collections = MagicMock()
        mock_collections.collections = []  # empty — no collections exist
        mock_client.get_collections.return_value = mock_collections

        writer = QdrantWriter()
        writer.client = mock_client

        writer.ensure_collection()
        mock_client.create_collection.assert_called_once()

    def test_ensure_collection_skips_when_exists(self):
        from spark.qdrant_writer import QdrantWriter

        mock_client = MagicMock()
        mock_collection = MagicMock()
        mock_collection.name = "academic_chunks"
        mock_collections = MagicMock()
        mock_collections.collections = [mock_collection]
        mock_client.get_collections.return_value = mock_collections

        writer = QdrantWriter()
        writer.client = mock_client

        writer.ensure_collection()
        mock_client.create_collection.assert_not_called()

    def test_upsert_chunks(self):
        from spark.qdrant_writer import QdrantWriter
        import numpy as np

        mock_client = MagicMock()
        writer = QdrantWriter()
        writer.client = mock_client

        chunks = [
            {
                "text": "Test chunk",
                "embedding": np.random.rand(384).tolist(),
                "chunk_index": 0,
                "token_count": 5,
                "source_label": "Page 1",
                "source_indices": [0],
                "metadata": {"page_number": 1},
            },
        ]

        count = writer.upsert_chunks(
            chunks=chunks,
            doc_id="test-doc",
            student_id="stu-001",
            course_id="cs-101",
        )

        assert count == 1
        mock_client.upsert.assert_called_once()

    def test_upsert_empty_chunks(self):
        from spark.qdrant_writer import QdrantWriter

        writer = QdrantWriter()
        writer.client = MagicMock()

        count = writer.upsert_chunks(chunks=[], doc_id="test-doc")
        assert count == 0


# ═══════════════════════════════════════════════════════════════
#  Full Pipeline Test (all external services mocked)
# ═══════════════════════════════════════════════════════════════


class TestProcessDocument:
    """End-to-end test of the process_document function."""

    def test_process_pdf_document(self):
        import numpy as np
        from spark.processing_job import process_document

        # Mock MinIO
        mock_minio = MagicMock()
        mock_response = MagicMock()
        mock_response.read.return_value = _make_pdf_bytes()
        mock_minio.get_object.return_value = mock_response

        # Mock Qdrant
        mock_qdrant = MagicMock()
        mock_qdrant.upsert_chunks.return_value = 5

        # Mock embedder to avoid loading the real model
        mock_model = MagicMock()
        mock_model.encode.return_value = np.random.rand(10, 384).astype(np.float32)

        message = {
            "doc_id": "test-123",
            "student_id": "stu-001",
            "course_id": "cs-101",
            "storage_path": "academic-documents/test/file.pdf",
            "file_type": ".pdf",
            "original_filename": "lecture.pdf",
        }

        with patch.dict("spark.embedder._MODEL_CACHE", {"test-model": mock_model}):
            result = process_document(
                message=message,
                minio_client=mock_minio,
                qdrant_writer=mock_qdrant,
                model_name="test-model",
            )

        assert result["doc_id"] == "test-123"
        assert result["status"] == "processed"
        assert result["error"] is None
        mock_qdrant.ensure_collection.assert_called_once()
        mock_qdrant.upsert_chunks.assert_called_once()

    def test_process_corrupted_file_returns_failed(self):
        from spark.processing_job import process_document

        # Mock MinIO to return garbage data
        mock_minio = MagicMock()
        mock_response = MagicMock()
        mock_response.read.return_value = b"corrupted garbage data"
        mock_minio.get_object.return_value = mock_response

        mock_qdrant = MagicMock()

        message = {
            "doc_id": "bad-doc",
            "student_id": "stu-001",
            "course_id": "cs-101",
            "storage_path": "academic-documents/test/bad.pdf",
            "file_type": ".pdf",
            "original_filename": "bad.pdf",
        }

        result = process_document(
            message=message,
            minio_client=mock_minio,
            qdrant_writer=mock_qdrant,
        )

        assert result["doc_id"] == "bad-doc"
        assert result["status"] == "failed"
        assert result["error"] is not None
        # Qdrant should NOT have been called
        mock_qdrant.upsert_chunks.assert_not_called()

    def test_process_text_document(self):
        import numpy as np
        from spark.processing_job import process_document

        mock_minio = MagicMock()
        mock_response = MagicMock()
        mock_response.read.return_value = _make_text_bytes()
        mock_minio.get_object.return_value = mock_response

        mock_qdrant = MagicMock()
        mock_qdrant.upsert_chunks.return_value = 3

        mock_model = MagicMock()
        mock_model.encode.return_value = np.random.rand(10, 384).astype(np.float32)

        message = {
            "doc_id": "txt-456",
            "student_id": "stu-002",
            "course_id": "ds-201",
            "storage_path": "academic-documents/test/notes.txt",
            "file_type": ".txt",
            "original_filename": "notes.txt",
        }

        with patch.dict("spark.embedder._MODEL_CACHE", {"test-model": mock_model}):
            result = process_document(
                message=message,
                minio_client=mock_minio,
                qdrant_writer=mock_qdrant,
                model_name="test-model",
            )

        assert result["status"] == "processed"
        assert result["error"] is None
