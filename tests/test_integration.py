"""
End-to-end integration tests — full pipeline flow.

Tests the complete path: upload -> parse -> chunk -> embed -> store -> retrieve -> answer.
External services (MinIO, Kafka, Qdrant, LLM) are mocked at the boundary,
but all internal pipeline code runs for real.
"""

import io
import json
import time
import uuid
import pytest
import numpy as np
from unittest.mock import MagicMock, patch
from fastapi.testclient import TestClient

from backend.app.main import app
from backend.app.services.document_service import document_service


# ===================================================================
#  Test file generators (real files, not stubs)
# ===================================================================

def make_pdf_bytes() -> bytes:
    import fitz
    doc = fitz.open()
    content = [
        "Machine learning is a branch of artificial intelligence that focuses on building systems that learn from data. Supervised learning uses labeled datasets to train predictive models. Common algorithms include linear regression, decision trees, and support vector machines.",
        "Deep learning is a subset of machine learning that uses neural networks with many layers. Convolutional neural networks are used for image recognition. Recurrent neural networks handle sequential data like text and time series.",
        "Reinforcement learning trains agents to make decisions by maximizing cumulative rewards. Key concepts include states, actions, policies, and value functions. Applications include game playing, robotics, and autonomous driving.",
    ]
    for text in content:
        page = doc.new_page()
        page.insert_text((72, 72), text, fontsize=11)
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def make_pptx_bytes() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    slides_data = [
        ("Introduction to NLP", "Natural language processing enables computers to understand human language. Key tasks include tokenization, named entity recognition, and sentiment analysis."),
        ("Text Classification", "Text classification assigns labels to documents. Approaches include bag-of-words, TF-IDF, and transformer-based models like BERT."),
        ("Language Models", "Language models predict the probability of word sequences. GPT and BERT are prominent examples. Pre-training on large corpora enables transfer learning."),
    ]
    for title, body in slides_data:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_text_bytes() -> bytes:
    return (
        "Computer vision allows machines to interpret visual data from the world.\n\n"
        "Object detection identifies and localizes objects within images using bounding boxes.\n\n"
        "Image segmentation assigns a class label to every pixel in an image.\n\n"
        "Convolutional neural networks are the backbone of modern computer vision systems.\n\n"
        "Transfer learning allows models trained on ImageNet to be fine-tuned for specific tasks.\n\n"
        "Data augmentation techniques like rotation, flipping, and cropping increase training data diversity."
    ).encode("utf-8")


# ===================================================================
#  Fixtures
# ===================================================================

@pytest.fixture(autouse=True)
def clear_docs():
    document_service._documents.clear()
    yield
    document_service._documents.clear()


@pytest.fixture
def client():
    return TestClient(app, raise_server_exceptions=False)


@pytest.fixture
def mock_upload_services():
    """Mock both MinIO and Kafka at the upload router level."""
    with patch("backend.app.routers.upload.minio_service") as mock_minio, \
         patch("backend.app.routers.upload.kafka_service") as mock_kafka:
        mock_minio.upload_file.return_value = "academic-documents/test/file.pdf"
        mock_kafka.publish.return_value = True
        yield mock_minio, mock_kafka


@pytest.fixture
def mock_embed_model():
    """Mock embedding model to return deterministic 384-dim vectors."""
    mock_model = MagicMock()

    def fake_encode(texts, **kwargs):
        if isinstance(texts, str):
            texts = [texts]
        vecs = []
        for t in texts:
            np.random.seed(hash(t[:20]) % (2**31))
            vecs.append(np.random.rand(384).astype(np.float32))
        return np.array(vecs)

    mock_model.encode.side_effect = fake_encode
    return mock_model


# ===================================================================
#  E2E Test: Upload -> Parse -> Chunk -> Embed -> Qdrant -> Ask
# ===================================================================

class TestEndToEndPipeline:
    """
    Full pipeline integration test.

    Simulates the complete lifecycle:
    1. Upload a PDF via the API
    2. Process it through the Spark pipeline (parse -> chunk -> embed)
    3. Store chunks in Qdrant (mocked)
    4. Query via /ask and verify the answer uses the right context
    """

    def test_pdf_upload_to_answer(self, client, mock_upload_services, mock_embed_model):
        """E2E: PDF upload -> processing -> retrieval -> grounded answer."""

        # -- Step 1: Upload --
        pdf_bytes = make_pdf_bytes()
        response = client.post(
            "/upload?student_id=stu-001&course_id=cs-101",
            files={"file": ("lecture.pdf", pdf_bytes, "application/pdf")},
        )
        assert response.status_code == 201
        upload_result = response.json()
        doc_id = upload_result["doc_id"]
        assert doc_id

        # Verify document is tracked
        status_resp = client.get(f"/documents/{doc_id}/status")
        assert status_resp.status_code == 200
        assert status_resp.json()["doc_id"] == doc_id

        # -- Step 2: Simulate Spark processing --
        from spark.parsers.pdf_parser import PDFParser
        from spark.chunker import chunk_sections
        from spark.embedder import embed_chunks_batch

        parser = PDFParser()
        sections = parser.parse(pdf_bytes, "lecture.pdf")
        assert len(sections) == 3, f"Expected 3 pages, got {len(sections)}"

        chunks = chunk_sections(sections, max_tokens=30, overlap_fraction=0.15)
        assert len(chunks) >= 1, f"Expected >=1 chunks, got {len(chunks)}"

        # Convert to dicts for embedding
        chunks_data = [
            {
                "text": c.text,
                "chunk_index": c.chunk_index,
                "token_count": c.token_count,
                "source_label": c.source_label,
                "source_indices": c.source_indices,
                "metadata": c.metadata,
            }
            for c in chunks
        ]

        # Embed with mocked model
        with patch.dict("spark.embedder._MODEL_CACHE", {"all-MiniLM-L6-v2": mock_embed_model}):
            embedded_chunks = embed_chunks_batch(chunks_data, model_name="all-MiniLM-L6-v2")

        assert all("embedding" in c for c in embedded_chunks)
        assert len(embedded_chunks[0]["embedding"]) == 384

        # -- Step 3: Simulate Qdrant storage --
        from spark.qdrant_writer import QdrantWriter
        mock_qdrant_client = MagicMock()
        writer = QdrantWriter()
        writer.client = mock_qdrant_client

        count = writer.upsert_chunks(
            chunks=embedded_chunks,
            doc_id=doc_id,
            student_id="stu-001",
            course_id="cs-101",
            original_filename="lecture.pdf",
        )
        assert count == len(embedded_chunks)
        mock_qdrant_client.upsert.assert_called()

        # -- Step 4: Query via /ask --
        search_results = [
            {
                "text": c["text"],
                "score": 0.95 - i * 0.05,
                "doc_id": doc_id,
                "original_filename": "lecture.pdf",
                "source_label": c["source_label"],
                "chunk_index": c["chunk_index"],
                "page_number": c["metadata"].get("page_number"),
                "slide_number": None,
                "slide_title": None,
            }
            for i, c in enumerate(embedded_chunks[:3])
        ]

        with patch("backend.app.routers.ask.embed_query", return_value=[0.1]*384), \
             patch("backend.app.routers.ask.qdrant_retrieval_service") as mock_qdrant, \
             patch("backend.app.routers.ask.llm_client") as mock_llm:

            mock_qdrant.search.return_value = search_results
            mock_llm.chat.return_value = (
                "Machine learning is a branch of AI that learns from data "
                "[Source: Page 1]. Deep learning uses neural networks with "
                "many layers [Source: Page 2]."
            )

            ask_resp = client.post("/ask", json={
                "question": "What is machine learning?",
                "doc_id": doc_id,
            })

        assert ask_resp.status_code == 200
        answer = ask_resp.json()
        assert "machine learning" in answer["answer"].lower()
        assert answer["chunks_used"] == 3
        assert len(answer["citations"]) == 3
        assert answer["citations"][0]["doc_id"] == doc_id

    def test_pptx_upload_to_mindmap(self, client, mock_upload_services, mock_embed_model):
        """E2E: PPTX upload -> processing -> mind map generation."""

        pptx_bytes = make_pptx_bytes()

        response = client.post(
            "/upload",
            files={"file": ("slides.pptx", pptx_bytes,
                            "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        assert response.status_code == 201
        doc_id = response.json()["doc_id"]

        # Process through Spark pipeline
        from spark.parsers.pptx_parser import PPTXParser
        from spark.chunker import chunk_pptx_sections

        parser = PPTXParser()
        sections = parser.parse(pptx_bytes, "slides.pptx")
        assert len(sections) >= 3

        chunks = chunk_pptx_sections(sections, max_tokens=100)
        assert len(chunks) >= 1

        # Generate mind map
        with patch("backend.app.routers.mindmap.embed_query", return_value=[0.1]*384), \
             patch("backend.app.routers.mindmap.qdrant_retrieval_service") as mock_qdrant, \
             patch("backend.app.routers.mindmap.llm_client") as mock_llm:

            mock_qdrant.search.return_value = [
                {"text": c.text, "score": 0.9, "doc_id": doc_id,
                 "original_filename": "slides.pptx", "source_label": c.source_label,
                 "chunk_index": 0, "page_number": None, "slide_number": None, "slide_title": None}
                for c in chunks[:3]
            ]
            mock_llm.chat_json.return_value = {
                "root": {
                    "label": "NLP Concepts",
                    "children": [
                        {"label": "Text Classification", "children": [
                            {"label": "Bag of Words", "children": []},
                            {"label": "BERT", "children": []},
                        ]},
                        {"label": "Language Models", "children": [
                            {"label": "GPT", "children": []},
                        ]},
                    ],
                }
            }

            resp = client.post("/mindmap", json={"doc_id": doc_id})

        assert resp.status_code == 200
        body = resp.json()
        assert "mindmap" in body["mermaid_syntax"]
        assert body["root"]["label"] == "NLP Concepts"

    def test_text_upload_to_quiz(self, client, mock_upload_services, mock_embed_model):
        """E2E: TXT upload -> processing -> quiz generation."""

        txt_bytes = make_text_bytes()

        response = client.post(
            "/upload",
            files={"file": ("notes.txt", txt_bytes, "text/plain")},
        )
        assert response.status_code == 201
        doc_id = response.json()["doc_id"]

        # Process
        from spark.parsers.text_parser import TextParser
        from spark.chunker import chunk_sections

        parser = TextParser()
        sections = parser.parse(txt_bytes, "notes.txt")
        chunks = chunk_sections(sections, max_tokens=80)
        assert len(chunks) >= 1

        # Generate quiz
        with patch("backend.app.routers.quiz.embed_query", return_value=[0.1]*384), \
             patch("backend.app.routers.quiz.qdrant_retrieval_service") as mock_qdrant, \
             patch("backend.app.routers.quiz.llm_client") as mock_llm:

            mock_qdrant.search.return_value = [
                {"text": c.text, "score": 0.9, "doc_id": doc_id,
                 "original_filename": "notes.txt", "source_label": c.source_label,
                 "chunk_index": 0, "page_number": None, "slide_number": None, "slide_title": None}
                for c in chunks[:5]
            ]
            mock_llm.chat_json.return_value = {
                "questions": [{
                    "question": "What enables machines to interpret visual data?",
                    "options": ["A) NLP", "B) Computer vision", "C) Reinforcement learning", "D) Clustering"],
                    "correct_answer": "B",
                    "explanation": "Computer vision processes visual data.",
                    "source_label": "Section 1",
                }]
            }

            resp = client.post("/quiz", json={"doc_id": doc_id, "num_questions": 1})

        assert resp.status_code == 200
        assert len(resp.json()["questions"]) == 1
        assert resp.json()["questions"][0]["correct_answer"] == "B"


# ===================================================================
#  Failure Handling Tests
# ===================================================================

class TestFailureHandling:
    """Demonstrate error paths and dead-letter routing."""

    def test_corrupted_pdf_fails_gracefully(self):
        """Corrupted PDF -> parse error -> documents.failed routing."""
        from spark.processing_job import process_document

        mock_minio = MagicMock()
        mock_response = MagicMock()
        mock_response.read.return_value = b"corrupted garbage not a real PDF"
        mock_minio.get_object.return_value = mock_response

        mock_qdrant = MagicMock()

        message = {
            "doc_id": "corrupt-001",
            "student_id": "stu-001",
            "course_id": "cs-101",
            "storage_path": "academic-documents/test/bad.pdf",
            "file_type": ".pdf",
            "original_filename": "corrupted.pdf",
        }

        result = process_document(
            message=message,
            minio_client=mock_minio,
            qdrant_writer=mock_qdrant,
        )

        assert result["status"] == "failed"
        assert result["error"] is not None
        assert result["chunk_count"] == 0
        mock_qdrant.upsert_chunks.assert_not_called()

    def test_corrupted_pptx_fails_gracefully(self):
        """Corrupted PPTX -> parse error -> documents.failed routing."""
        from spark.processing_job import process_document

        mock_minio = MagicMock()
        mock_response = MagicMock()
        mock_response.read.return_value = b"not a valid pptx file at all"
        mock_minio.get_object.return_value = mock_response

        result = process_document(
            message={
                "doc_id": "corrupt-002",
                "student_id": "stu-001",
                "course_id": "cs-101",
                "storage_path": "academic-documents/test/bad.pptx",
                "file_type": ".pptx",
                "original_filename": "corrupted.pptx",
            },
            minio_client=mock_minio,
            qdrant_writer=MagicMock(),
        )

        assert result["status"] == "failed"
        assert result["error"] is not None

    def test_empty_text_file_fails(self):
        """Empty text file -> zero chunks -> failure."""
        from spark.processing_job import process_document

        mock_minio = MagicMock()
        mock_response = MagicMock()
        mock_response.read.return_value = b""
        mock_minio.get_object.return_value = mock_response

        result = process_document(
            message={
                "doc_id": "empty-001",
                "student_id": "stu-001",
                "course_id": "cs-101",
                "storage_path": "academic-documents/test/empty.txt",
                "file_type": ".txt",
                "original_filename": "empty.txt",
            },
            minio_client=mock_minio,
            qdrant_writer=MagicMock(),
        )

        assert result["status"] == "failed"

    def test_upload_unsupported_file_type(self, client, mock_upload_services):
        """Unsupported file extension rejected at API layer."""
        response = client.post(
            "/upload",
            files={"file": ("malware.exe", b"bad data", "application/octet-stream")},
        )
        assert response.status_code == 400

    def test_ask_with_empty_question_rejected(self, client):
        """Too-short questions are rejected by validation."""
        response = client.post("/ask", json={"question": "ab"})
        assert response.status_code == 422

    def test_quiz_with_no_content_returns_404(self, client):
        """Quiz on nonexistent doc returns 404."""
        with patch("backend.app.routers.quiz.embed_query", return_value=[0.1]*384), \
             patch("backend.app.routers.quiz.qdrant_retrieval_service") as mock_qdrant:
            mock_qdrant.search.return_value = []
            response = client.post("/quiz", json={"doc_id": "nonexistent"})
        assert response.status_code == 404


# ===================================================================
#  Performance Benchmarks
# ===================================================================

class TestPerformanceBenchmarks:
    """
    Measure component-level latencies and throughput.

    These are not correctness tests -- they measure performance
    and print results for documentation.
    """

    def test_pdf_parse_latency(self):
        """Benchmark: PDF parsing speed."""
        from spark.parsers.pdf_parser import PDFParser

        pdf_bytes = make_pdf_bytes()
        parser = PDFParser()

        # Warm up
        parser.parse(pdf_bytes, "warm.pdf")

        times = []
        for i in range(10):
            start = time.perf_counter()
            parser.parse(pdf_bytes, f"bench_{i}.pdf")
            times.append(time.perf_counter() - start)

        avg_ms = (sum(times) / len(times)) * 1000
        p95_ms = sorted(times)[int(0.95 * len(times))] * 1000
        print(f"\n[BENCH] PDF Parse: avg={avg_ms:.1f}ms  p95={p95_ms:.1f}ms  (3 pages)")
        assert avg_ms < 500, f"PDF parse too slow: {avg_ms:.1f}ms"

    def test_pptx_parse_latency(self):
        """Benchmark: PPTX parsing speed."""
        from spark.parsers.pptx_parser import PPTXParser

        pptx_bytes = make_pptx_bytes()
        parser = PPTXParser()
        parser.parse(pptx_bytes, "warm.pptx")

        times = []
        for i in range(10):
            start = time.perf_counter()
            parser.parse(pptx_bytes, f"bench_{i}.pptx")
            times.append(time.perf_counter() - start)

        avg_ms = (sum(times) / len(times)) * 1000
        p95_ms = sorted(times)[int(0.95 * len(times))] * 1000
        print(f"\n[BENCH] PPTX Parse: avg={avg_ms:.1f}ms  p95={p95_ms:.1f}ms  (3 slides)")
        assert avg_ms < 500, f"PPTX parse too slow: {avg_ms:.1f}ms"

    def test_chunking_latency(self):
        """Benchmark: chunking speed on varying text sizes."""
        from spark.parsers.base import ParsedSection
        from spark.chunker import chunk_sections

        sections = [
            ParsedSection(
                label=f"Page {i+1}",
                index=i,
                text=(
                    f"Section {i+1} discusses important concepts. "
                    "Machine learning algorithms learn patterns from data. "
                    "These patterns enable predictions on unseen examples. "
                    "Training requires careful hyperparameter tuning. "
                    "Cross-validation helps estimate generalization performance."
                ) * 3,
            )
            for i in range(50)
        ]

        start = time.perf_counter()
        chunks = chunk_sections(sections, max_tokens=100, overlap_fraction=0.15)
        elapsed_ms = (time.perf_counter() - start) * 1000

        print(f"\n[BENCH] Chunking: {elapsed_ms:.1f}ms for 50 sections -> {len(chunks)} chunks")
        assert elapsed_ms < 2000, f"Chunking too slow: {elapsed_ms:.1f}ms"

    def test_embedding_throughput(self, mock_embed_model):
        """Benchmark: embedding throughput (mocked model)."""
        from spark.embedder import embed_chunks_batch

        chunks_data = [
            {"text": f"Sample text for chunk number {i} about machine learning concepts."}
            for i in range(100)
        ]

        with patch.dict("spark.embedder._MODEL_CACHE", {"all-MiniLM-L6-v2": mock_embed_model}):
            start = time.perf_counter()
            result = embed_chunks_batch(chunks_data, model_name="all-MiniLM-L6-v2")
            elapsed_ms = (time.perf_counter() - start) * 1000

        throughput = len(chunks_data) / (elapsed_ms / 1000) if elapsed_ms > 0 else 0
        print(f"\n[BENCH] Embedding: {elapsed_ms:.1f}ms for {len(chunks_data)} chunks ({throughput:.0f} chunks/sec)")
        assert len(result) == 100

    def test_qdrant_upsert_throughput(self):
        """Benchmark: Qdrant upsert throughput (mocked client)."""
        from spark.qdrant_writer import QdrantWriter

        mock_client = MagicMock()
        writer = QdrantWriter()
        writer.client = mock_client

        chunks = [
            {
                "text": f"Chunk {i} text about AI.",
                "embedding": np.random.rand(384).tolist(),
                "chunk_index": i,
                "token_count": 20,
                "source_label": f"Page {i+1}",
                "source_indices": [i],
                "metadata": {"page_number": i + 1},
            }
            for i in range(200)
        ]

        start = time.perf_counter()
        count = writer.upsert_chunks(
            chunks=chunks, doc_id="bench-doc",
            student_id="stu-001", batch_size=50,
        )
        elapsed_ms = (time.perf_counter() - start) * 1000

        throughput = count / (elapsed_ms / 1000) if elapsed_ms > 0 else 0
        print(f"\n[BENCH] Qdrant Upsert: {elapsed_ms:.1f}ms for {count} chunks ({throughput:.0f} chunks/sec)")
        assert count == 200
        assert mock_client.upsert.call_count == 4  # 200 / batch_size=50

    def test_full_pipeline_latency(self, mock_embed_model):
        """Benchmark: complete pipeline (parse+chunk+embed+upsert) per document."""
        from spark.processing_job import process_document

        mock_minio = MagicMock()
        mock_response = MagicMock()
        mock_response.read.return_value = make_pdf_bytes()
        mock_minio.get_object.return_value = mock_response

        mock_qdrant = MagicMock()
        mock_qdrant.upsert_chunks.return_value = 10

        message = {
            "doc_id": "perf-test",
            "student_id": "stu-001",
            "course_id": "cs-101",
            "storage_path": "academic-documents/test/lecture.pdf",
            "file_type": ".pdf",
            "original_filename": "lecture.pdf",
        }

        with patch.dict("spark.embedder._MODEL_CACHE", {"all-MiniLM-L6-v2": mock_embed_model}):
            start = time.perf_counter()
            result = process_document(
                message=message,
                minio_client=mock_minio,
                qdrant_writer=mock_qdrant,
                model_name="all-MiniLM-L6-v2",
            )
            elapsed_ms = (time.perf_counter() - start) * 1000

        assert result["status"] == "processed"
        print(f"\n[BENCH] Full Pipeline: {elapsed_ms:.1f}ms per document (parse+chunk+embed+upsert)")
        assert elapsed_ms < 5000

    def test_api_ask_latency(self, client):
        """Benchmark: /ask endpoint latency (all services mocked)."""
        with patch("backend.app.routers.ask.embed_query", return_value=[0.1]*384), \
             patch("backend.app.routers.ask.qdrant_retrieval_service") as mock_qdrant, \
             patch("backend.app.routers.ask.llm_client") as mock_llm:

            mock_qdrant.search.return_value = [
                {"text": "ML is AI.", "score": 0.9, "doc_id": "d1",
                 "original_filename": "f.pdf", "source_label": "Page 1",
                 "chunk_index": 0, "page_number": 1, "slide_number": None, "slide_title": None}
            ]
            mock_llm.chat.return_value = "Answer text."

            times = []
            for _ in range(20):
                start = time.perf_counter()
                resp = client.post("/ask", json={"question": "What is ML?"})
                times.append(time.perf_counter() - start)
                assert resp.status_code == 200

            avg_ms = (sum(times) / len(times)) * 1000
            p50_ms = sorted(times)[len(times)//2] * 1000
            p95_ms = sorted(times)[int(0.95 * len(times))] * 1000

        print(f"\n[BENCH] /ask Latency: p50={p50_ms:.1f}ms  avg={avg_ms:.1f}ms  p95={p95_ms:.1f}ms  (n=20)")
